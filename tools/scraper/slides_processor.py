import os
from typing import Dict, Any, List
from pptx import Presentation as PPTXPresentation

class SlidesProcessor:
    """Extracts text and structure from Google Slides (API) and PowerPoint (.pptx) files."""

    def process_google_slides(self, presentation: Dict[str, Any]) -> Dict[str, Any]:
        """Converts a Google Slides API structure into a Dict of {markdown: str, images: {id: url}}."""
        title = presentation.get('title', 'Untitled Presentation')
        markdown = [f"# Slides: {title}\n"]
        images = {}
        
        slides = presentation.get('slides', [])
        for i, slide in enumerate(slides):
            markdown.append(f"## Slide {i + 1}")
            
            elements = slide.get('pageElements', [])
            elements.sort(key=lambda x: x.get('transform', {}).get('translateY', 0))
            
            for element in elements:
                text = self._extract_text_from_google_element(element)
                if text.strip():
                    markdown.append(text.strip())
                
                # Image detection
                if 'image' in element:
                    obj_id = element.get('objectId')
                    url = element['image'].get('contentUrl')
                    if obj_id and url:
                        images[obj_id] = url
                        markdown.append(f"\n\n[IMAGE:{obj_id}]\n\n")
            
            markdown.append("\n---\n")
            
        return {"markdown": "\n\n".join(markdown), "images": images}

    def process_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """Extracts text and embedded images from a local .pptx file."""
        basename = os.path.basename(pptx_path)
        markdown = [f"# PowerPoint: {basename}\n"]
        images = {} # Will store mapping of internal image name to binary index if needed
        try:
            prs = PPTXPresentation(pptx_path)
            for i, slide in enumerate(prs.slides):
                markdown.append(f"## Slide {i + 1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        markdown.append(shape.text.strip())
                    if shape.has_table:
                        rows = []
                        for row in shape.table.rows:
                            cells = [cell.text_frame.text.strip().replace('\n', ' ') for cell in row.cells]
                            rows.append("| " + " | ".join(cells) + " |")
                        markdown.append("\n".join(rows))
                    
                    # PPTX Image detection
                    if shape.shape_type == 13: # 13 is PICTURE
                        # For PPTX, we might want to extract the image bytes directly
                        obj_id = f"pptx_img_{i}_{getattr(shape, 'shape_id', 'unknown')}"
                        markdown.append(f"\n\n[IMAGE:{obj_id}]\n\n")
                        # This id can be used by main.py to pull shape.image.blob
                        images[obj_id] = shape 
                        
                markdown.append("\n---\n")
            return {"markdown": "\n\n".join(markdown), "images": images}
        except Exception as e:
            print(f"Error processing PPTX {pptx_path}: {e}")
            return {"markdown": f"Error: Could not process {pptx_path}", "images": {}}

    def _extract_text_from_google_element(self, element: Dict[str, Any]) -> str:
        """Extracts text from various Google Slides API elements."""
        text_parts = []
        
        # Handle regular shapes (text boxes, titles)
        if 'shape' in element:
            text_content = element['shape'].get('text', {})
            for text_element in text_content.get('textElements', []):
                if 'textRun' in text_element:
                    text_parts.append(text_element['textRun'].get('content', ''))
        
        # Handle tables
        elif 'table' in element:
            table_rows = element['table'].get('tableRows', [])
            for row in table_rows:
                row_parts = []
                for cell in row.get('tableCells', []):
                    cell_text = ""
                    cell_content = cell.get('text', {})
                    for text_element in cell_content.get('textElements', []):
                        if 'textRun' in text_element:
                            cell_text += text_element['textRun'].get('content', '')
                    row_parts.append(cell_text.strip().replace('\n', ' '))
                text_parts.append("| " + " | ".join(row_parts) + " |")
                
        return "".join(text_parts)

# (os import moved to top)

if __name__ == "__main__":
    # Small test dummy
    processor = SlidesProcessor()
    dummy_pres = {
        "title": "Test Presentation",
        "slides": [
            {
                "pageElements": [
                    {"shape": {"text": {"textElements": [{"textRun": {"content": "Hello World\n"}}]}}}
                ]
            }
        ]
    }
    print(processor.process_google_slides(dummy_pres))
